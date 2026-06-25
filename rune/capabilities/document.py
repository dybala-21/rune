"""Read and write office documents (xlsx, pptx, docx, pdf, csv, html).

Renders structured blocks/sheets into files and extracts text back out, using
pure-python libraries. No code execution or network.
"""

from __future__ import annotations

import csv as _csv
import html as _html
import os
from pathlib import Path
from typing import Literal

from pydantic import BaseModel, Field

from rune.capabilities.registry import CapabilityRegistry
from rune.capabilities.types import CapabilityDefinition
from rune.safety.guardian import get_guardian
from rune.types import CapabilityResult, Domain, RiskLevel
from rune.utils.logger import get_logger

log = get_logger(__name__)

# Parameter schema

DocFormat = Literal["xlsx", "pptx", "docx", "pdf", "csv", "html"]


class DocBlock(BaseModel):
    """One content block. Rendered per-format by each renderer."""

    type: Literal["heading", "paragraph", "bullets", "table", "page_break"]
    text: str = Field(default="", description="Text for heading/paragraph")
    level: int = Field(default=1, description="Heading level 1-3")
    items: list[str] = Field(default_factory=list, description="Bullet items")
    rows: list[list[str | int | float]] = Field(
        default_factory=list,
        description="Table rows; the first row is treated as the header",
    )


class DocSheet(BaseModel):
    """One spreadsheet tab (xlsx/csv). First row is treated as the header."""

    name: str = "Sheet1"
    rows: list[list[str | int | float]] = Field(default_factory=list)


class DocumentCreateParams(BaseModel):
    path: str = Field(description="Output file path; extension should match format")
    format: DocFormat = Field(description="xlsx, pptx, docx, pdf, csv, or html")
    title: str = Field(default="", description="Document or first-slide title")
    blocks: list[DocBlock] = Field(
        default_factory=list,
        description="Content blocks for pptx/docx/pdf/html",
    )
    sheets: list[DocSheet] = Field(
        default_factory=list, description="Spreadsheet tabs for xlsx/csv"
    )


# Helpers

def _sheets_from_blocks(blocks: list[DocBlock]) -> list[DocSheet]:
    """Fall back to table blocks when no explicit sheets are given."""
    sheets: list[DocSheet] = []
    for i, b in enumerate(blocks):
        if b.type == "table" and b.rows:
            sheets.append(DocSheet(name=f"Sheet{i + 1}", rows=b.rows))
    return sheets


def _missing_dep_error(fmt: str, mod: str) -> CapabilityResult:
    return CapabilityResult(
        success=False,
        error=(
            f"{fmt} support needs the '{mod}' package. It ships in RUNE's "
            f"dependencies; run `uv sync` to install it."
        ),
    )


# Renderers

def _render_xlsx(p: Path, params: DocumentCreateParams) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    wb.remove(wb.active)
    sheets = params.sheets or _sheets_from_blocks(params.blocks)
    if not sheets:
        sheets = [DocSheet(name="Sheet1", rows=[["(empty)"]])]
    for s in sheets:
        ws = wb.create_sheet(title=(s.name or "Sheet")[:31])
        for row in s.rows:
            ws.append(list(row))
        if s.rows:
            for cell in ws[1]:
                cell.font = Font(bold=True)
    wb.save(str(p))


def _render_csv(p: Path, params: DocumentCreateParams) -> None:
    sheets = params.sheets or _sheets_from_blocks(params.blocks)
    rows = sheets[0].rows if sheets else [["(empty)"]]
    with open(p, "w", newline="", encoding="utf-8") as f:
        _csv.writer(f).writerows(rows)


def _render_docx(p: Path, params: DocumentCreateParams) -> None:
    from docx import Document

    doc = Document()
    if params.title:
        doc.add_heading(params.title, level=0)
    for b in params.blocks:
        if b.type == "heading":
            doc.add_heading(b.text, level=max(1, min(b.level, 4)))
        elif b.type == "paragraph":
            doc.add_paragraph(b.text)
        elif b.type == "bullets":
            for item in b.items:
                doc.add_paragraph(item, style="List Bullet")
        elif b.type == "table" and b.rows:
            cols = max(len(r) for r in b.rows)
            t = doc.add_table(rows=0, cols=cols)
            t.style = "Light Grid Accent 1"
            for r in b.rows:
                cells = t.add_row().cells
                for j, val in enumerate(r):
                    cells[j].text = str(val)
        elif b.type == "page_break":
            doc.add_page_break()
    doc.save(str(p))


def _render_pptx(p: Path, params: DocumentCreateParams) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    if params.title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = params.title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = ""

    current_body = None

    def _new_content_slide(heading: str) -> None:
        nonlocal current_body
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = heading
        current_body = slide.placeholders[1].text_frame
        current_body.clear()

    for b in params.blocks:
        if b.type == "heading":
            _new_content_slide(b.text)
        elif b.type in ("paragraph", "bullets"):
            if current_body is None:
                _new_content_slide(params.title or "Slide")
            texts = b.items if b.type == "bullets" else [b.text]
            for txt in texts:
                para = current_body.add_paragraph()
                para.text = txt
        elif b.type == "table" and b.rows:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            rows, cols = len(b.rows), max(len(r) for r in b.rows)
            shape = slide.shapes.add_table(
                rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * rows)
            )
            for i, r in enumerate(b.rows):
                for j in range(cols):
                    shape.table.cell(i, j).text = str(r[j]) if j < len(r) else ""
            current_body = None
    if not prs.slides:
        prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(p))


def _render_pdf(p: Path, params: DocumentCreateParams) -> None:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    # Core fonts are latin-1 only.
    def _s(text: str) -> str:
        return text.encode("latin-1", "replace").decode("latin-1")

    # Return to the left margin so the next multi_cell has room.
    def _mcell(h: float, text: str) -> None:
        pdf.multi_cell(0, h, text, new_x="LMARGIN", new_y="NEXT")

    if params.title:
        pdf.set_font("Helvetica", "B", 20)
        _mcell(10, _s(params.title))
        pdf.ln(2)
    for b in params.blocks:
        if b.type == "heading":
            pdf.set_font("Helvetica", "B", 16 - min(b.level, 3) * 2)
            _mcell(8, _s(b.text))
        elif b.type == "paragraph":
            pdf.set_font("Helvetica", "", 11)
            _mcell(6, _s(b.text))
        elif b.type == "bullets":
            pdf.set_font("Helvetica", "", 11)
            for item in b.items:
                _mcell(6, _s(f"- {item}"))
        elif b.type == "table" and b.rows:
            pdf.set_font("Helvetica", "", 10)
            cols = max(len(r) for r in b.rows)
            w = pdf.epw / max(cols, 1)  # effective page width / columns
            for r in b.rows:
                for j in range(cols):
                    cell = _s(str(r[j]) if j < len(r) else "")
                    pdf.cell(w, 7, cell[:24], border=1)
                pdf.ln(7)
        elif b.type == "page_break":
            pdf.add_page()
        pdf.ln(1)
    pdf.output(str(p))


def _render_html(p: Path, params: DocumentCreateParams) -> None:
    parts = ["<!doctype html><html><head><meta charset='utf-8'>"]
    if params.title:
        parts.append(f"<title>{_html.escape(params.title)}</title>")
    parts.append("</head><body>")
    if params.title:
        parts.append(f"<h1>{_html.escape(params.title)}</h1>")
    for b in params.blocks:
        if b.type == "heading":
            lvl = max(1, min(b.level + 1, 6))
            parts.append(f"<h{lvl}>{_html.escape(b.text)}</h{lvl}>")
        elif b.type == "paragraph":
            parts.append(f"<p>{_html.escape(b.text)}</p>")
        elif b.type == "bullets":
            parts.append("<ul>")
            parts.extend(f"<li>{_html.escape(i)}</li>" for i in b.items)
            parts.append("</ul>")
        elif b.type == "table" and b.rows:
            parts.append("<table border='1' cellspacing='0' cellpadding='4'>")
            for ri, r in enumerate(b.rows):
                tag = "th" if ri == 0 else "td"
                cells = "".join(f"<{tag}>{_html.escape(str(c))}</{tag}>" for c in r)
                parts.append(f"<tr>{cells}</tr>")
            parts.append("</table>")
        elif b.type == "page_break":
            parts.append("<hr>")
    parts.append("</body></html>")
    p.write_text("".join(parts), encoding="utf-8")


_RENDERERS = {
    "xlsx": (_render_xlsx, "openpyxl"),
    "csv": (_render_csv, ""),
    "docx": (_render_docx, "python-docx"),
    "pptx": (_render_pptx, "python-pptx"),
    "pdf": (_render_pdf, "fpdf2"),
    "html": (_render_html, ""),
}


async def document_create(params: DocumentCreateParams) -> CapabilityResult:
    """Render a structured spec into a real document file."""
    if not params.path or not params.path.strip():
        return CapabilityResult(success=False, error="Empty file path")

    guardian = get_guardian()
    validation = guardian.validate_file_path(params.path)
    if not validation.allowed:
        return CapabilityResult(success=False, error=validation.reason)

    file_path = Path(os.path.expanduser(params.path)).resolve()
    home = os.environ.get("HOME", str(Path.home()))
    if str(file_path) in ("/", home) or len(file_path.parts) < 3:
        return CapabilityResult(
            success=False,
            error=f"BLOCKED: refusing to write to critical path: {file_path}",
        )

    renderer_entry = _RENDERERS.get(params.format)
    if renderer_entry is None:
        return CapabilityResult(
            success=False, error=f"Unsupported format: {params.format}"
        )
    renderer, mod = renderer_entry

    file_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        renderer(file_path, params)
    except ImportError:
        return _missing_dep_error(params.format, mod)
    except Exception as exc:
        log.debug("document_create_failed", format=params.format, error=str(exc))
        return CapabilityResult(
            success=False, error=f"Failed to render {params.format}: {exc}"
        )

    size = file_path.stat().st_size if file_path.is_file() else 0
    return CapabilityResult(
        success=True,
        output=f"Created {params.format} ({size} bytes): {params.path}",
        metadata={"path": str(file_path), "format": params.format, "bytes": size},
    )


# Reading: extract text/tables from existing documents

class DocumentReadParams(BaseModel):
    path: str = Field(description="Path to the document to read")
    max_chars: int = Field(
        default=20_000,
        description="Cap on extracted text returned (longer content is truncated)",
    )


def _read_xlsx(p: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(p), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(cells):
                out.append("\t".join(cells))
    return "\n".join(out)


def _read_docx(p: Path) -> str:
    from docx import Document

    doc = Document(str(p))
    out = [para.text for para in doc.paragraphs if para.text]
    for t in doc.tables:
        for row in t.rows:
            out.append("\t".join(c.text for c in row.cells))
    return "\n".join(out)


def _read_pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                out.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    out.append("\t".join(c.text for c in row.cells))
    return "\n".join(out)


def _read_pdf(p: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(p))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _read_text(p: Path) -> str:
    return p.read_text(encoding="utf-8", errors="replace")


_READERS = {
    "xlsx": (_read_xlsx, "openpyxl"),
    "docx": (_read_docx, "python-docx"),
    "pptx": (_read_pptx, "python-pptx"),
    "pdf": (_read_pdf, "pypdf"),
    "csv": (_read_text, ""),
    "html": (_read_text, ""),
    "txt": (_read_text, ""),
    "md": (_read_text, ""),
}


async def document_read(params: DocumentReadParams) -> CapabilityResult:
    """Extract text and tables from an existing document file."""
    if not params.path or not params.path.strip():
        return CapabilityResult(success=False, error="Empty file path")

    guardian = get_guardian()
    validation = guardian.validate_file_path(params.path)
    if not validation.allowed:
        return CapabilityResult(success=False, error=validation.reason)

    file_path = Path(os.path.expanduser(params.path)).resolve()
    if not file_path.is_file():
        return CapabilityResult(success=False, error=f"File not found: {params.path}")

    ext = file_path.suffix.lower().lstrip(".")
    reader_entry = _READERS.get(ext)
    if reader_entry is None:
        return CapabilityResult(
            success=False,
            error=f"Unsupported document type '.{ext}'. Supported: {sorted(_READERS)}",
        )
    reader, mod = reader_entry

    try:
        text = reader(file_path)
    except ImportError:
        return _missing_dep_error(ext, mod)
    except Exception as exc:
        log.debug("document_read_failed", ext=ext, error=str(exc))
        return CapabilityResult(success=False, error=f"Failed to read {ext}: {exc}")

    truncated = len(text) > params.max_chars
    if truncated:
        text = text[: params.max_chars] + f"\n... [truncated, {len(text)} total chars]"
    return CapabilityResult(
        success=True,
        output=text,
        metadata={"path": str(file_path), "format": ext, "truncated": truncated},
    )


def register_document_capability(registry: CapabilityRegistry) -> None:
    """Register the document_create and document_read capabilities."""
    registry.register(CapabilityDefinition(
        name="document_create",
        description=(
            "Create a real document file (xlsx, pptx, docx, pdf, csv, html) from "
            "structured content. Use 'sheets' for xlsx/csv and 'blocks' "
            "(heading/paragraph/bullets/table/page_break) for pptx/docx/pdf/html."
        ),
        domain=Domain.FILE,
        risk_level=RiskLevel.MEDIUM,
        group="file",
        parameters_model=DocumentCreateParams,
        execute=document_create,
    ))
    registry.register(CapabilityDefinition(
        name="document_read",
        description=(
            "Extract text and tables from an existing document "
            "(xlsx, docx, pptx, pdf, csv, html, txt, md). Use this to read or "
            "summarize spreadsheets, slide decks, Word, and PDF files."
        ),
        domain=Domain.FILE,
        risk_level=RiskLevel.LOW,
        group="file",
        parameters_model=DocumentReadParams,
        execute=document_read,
    ))
